import os
import json
import re
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configure Gemini API
client = genai.Client()

def parse_script_to_slides(script_text: str) -> list:
    """
    Use Gemini to parse a video script into a structured JSON array for slides.
    Each item in the array represents a slide: {"title": "...", "bullets": ["...", "..."]}
    """
    system_prompt = """
You are an expert Presentation Creator. 
Your task is to read the provided video script or text and summarize it into a series of presentation slides.
Return the output STRICTLY as a JSON array of objects.
Each object must have exactly two keys:
- "title": A short, catchy title for the slide (string)
- "bullets": A list of 1 to 4 short bullet points summarizing the content (list of strings)

Do not output any markdown formatting other than the JSON block itself. Do not include any explanations.
"""
    
    response = client.models.generate_content(
        model='gemini-2.5-flash',
        contents=f"{system_prompt}\n\n[Script]\n{script_text}",
    )
    
    # Parse the response to extract JSON
    text = response.text
    # Try to find JSON block
    match = re.search(r'```(?:json)?\n(.*?)\n```', text, re.DOTALL)
    if match:
        json_str = match.group(1)
    else:
        json_str = text.strip()
        
    try:
        slide_data = json.loads(json_str)
        return slide_data
    except json.JSONDecodeError as e:
        print(f"Failed to parse JSON from LLM: {text}")
        raise e

def generate_pptx(slide_data: list, output_path: str):
    """
    Generate a PowerPoint presentation using python-pptx.
    Applies a basic dark theme (black background, white text) since no template is provided.
    """
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # slide_layouts[1] is typically Title and Content
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_info in slide_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set Dark Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30) # Dark Gray / Black
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        # Title
        title_shape.text = slide_info.get("title", "No Title")
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.name = "Malgun Gothic"
        
        # Bullets
        tf = body_shape.text_frame
        tf.clear() # Clear default text
        
        bullets = slide_info.get("bullets", [])
        for i, bullet_text in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = bullet_text
            p.level = 0
            
            # Formatting bullet text
            for run in p.runs:
                run.font.color.rgb = RGBColor(220, 220, 220)
                run.font.size = Pt(28)
                run.font.name = "Malgun Gothic"
    
    # Create directory if not exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    # For quick testing
    sample_script = "오늘은 VibeCoding 파이프라인 구축 완료 보고를 드립니다. 첫째, Claude CLI 대화형 래핑을 걷어내고 파이썬 내장 API 직접 호출 방식으로 아키텍처를 전면 개편했습니다. 둘째, 단순 색상이나 레이아웃 변경은 LLM을 거치지 않도록 정규식(Regex) 화이트리스트를 구현하여 비용을 100% 절감했습니다. 마지막으로 이 모든 파이프라인에 5분 타임아웃 강제 종료 기능을 탑재하여 무한 대기를 원천 차단했습니다. 만복이에게도 앞으로 대화형 툴의 백그라운드 래핑은 무조건 피하라고 전파 완료했습니다."
    
    print("Generating slide data from script...")
    data = parse_script_to_slides(sample_script)
    
    print(f"Generated {len(data)} slides. Creating PPTX...")
    out_path = os.path.join(os.path.dirname(__file__), "output", "test_presentation.pptx")
    generate_pptx(data, out_path)
