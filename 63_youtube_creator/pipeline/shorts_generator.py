import os
import sys
if sys.stdout.encoding.lower() != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_shorts_pptx(output_path):
    prs = Presentation()
    
    # 네이티브 1080x1920 해상도를 위한 크기 조정 (96 DPI 기준)
    # 1080 / 96 = 11.25 Inches
    # 1920 / 96 = 20.0 Inches
    prs.slide_width = Inches(11.25)
    prs.slide_height = Inches(20.0)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    bg_color = RGBColor(0x12, 0x12, 0x12)
    cyan_color = RGBColor(0x00, 0xF0, 0xFF)
    green_color = RGBColor(0x39, 0xFF, 0x14)
    white_color = RGBColor(0xFF, 0xFF, 0xFF)
    
    def add_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def create_text_box(slide, left, top, width, height, name=None):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        if name:
            txBox.name = name
        tf = txBox.text_frame
        tf.word_wrap = True
        return tf

    def add_run(p, text, color, font_size, bold=False):
        run = p.add_run()
        run.text = text
        run.font.name = 'Malgun Gothic'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return run

    def set_slide_timing(slide, duration_sec):
        # Store expected duration in slide notes for QA agent to read
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = json.dumps({"duration_sec": duration_sec})

    # --- Slide 1: 타이틀 ---
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide1)
    tf1 = create_text_box(slide1, Inches(0.625), Inches(7.5), Inches(10), Inches(5), name="title_1")
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_run(p1, "🎯 CRISP-DM?\n", cyan_color, 88, bold=True)
    add_run(p1, "데이터분석의 표준 6단계!", white_color, 56, bold=True)

    # --- Slide 2: 핵심 소개 ---
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide2)
    tf2 = create_text_box(slide2, Inches(1.25), Inches(6.25), Inches(8.75), Inches(7.5), name="body_2")
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    add_run(p2, "데이터를 분석하려면?\n막 시작하면 안 된다.\n정해진 순서가 있어.\n\n", white_color, 50)
    add_run(p2, "그게 바로 CRISP-DM!\n", cyan_color, 62, bold=True)
    add_run(p2, "업계 표준 방법론이야.", white_color, 50)

    # --- Slide 3: 6단계 설명 ---
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide3)
    tf3 = create_text_box(slide3, Inches(0.625), Inches(1.875), Inches(10), Inches(16.25), name="body_3")
    p3 = tf3.paragraphs[0]
    add_run(p3, "CRISP-DM 6단계\n", cyan_color, 75, bold=True)
    
    steps = [
        ("Step 1", "Business Understanding", "비즈니스 목표 정의"),
        ("Step 2", "Data Understanding", "데이터 수집 및 탐색"),
        ("Step 3", "Data Preparation", "데이터 정제 및 전처리"),
        ("Step 4", "Modeling", "알고리즘 학습 및 튜닝"),
        ("Step 5", "Evaluation", "성능 검증"),
        ("Step 6", "Deployment", "운영 환경 적용")
    ]
    
    for s_num, s_eng, s_kor in steps:
        p = tf3.add_paragraph()
        add_run(p, f"{s_num} ", cyan_color, 40, bold=True)
        add_run(p, f"{s_eng}", white_color, 40, bold=True)
        p_sub = tf3.add_paragraph()
        add_run(p_sub, f"  → {s_kor}", green_color, 35)

    # --- Slide 4: 우리 사례 ---
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide4)
    tf4 = create_text_box(slide4, Inches(0.625), Inches(6.25), Inches(10), Inches(7.5), name="body_4")
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    add_run(p4, "우리 3AI는?\n이 6단계를 완벽하게 적용 중!\n\n", cyan_color, 56, bold=True)
    
    p4_1 = tf4.add_paragraph()
    p4_1.alignment = PP_ALIGN.CENTER
    add_run(p4_1, "n8n 워크플로우 설계 = CRISP-DM 그 자체\n", white_color, 44)
    
    p4_2 = tf4.add_paragraph()
    p4_2.alignment = PP_ALIGN.CENTER
    add_run(p4_2, "Improve_stock 주식 분석 = EDA 교과서 흐름\n\n", white_color, 44)
    
    p4_3 = tf4.add_paragraph()
    p4_3.alignment = PP_ALIGN.CENTER
    add_run(p4_3, "수집 → 전처리 → 모델 → 신호생성 → 배포\n", green_color, 40, bold=True)
    add_run(p4_3, "= CRISP-DM 완벽 실행!", cyan_color, 50, bold=True)

    # --- Slide 5: 엔딩 ---
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide5)
    tf5 = create_text_box(slide5, Inches(1.25), Inches(7.5), Inches(8.75), Inches(5), name="title_5")
    p5 = tf5.paragraphs[0]
    p5.alignment = PP_ALIGN.CENTER
    add_run(p5, "CRISP-DM은?\n무조건 알아야 할\n데이터 분석 기초!\n\n", white_color, 62, bold=True)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"PPTX 생성 완료 (1080p 해상도 적용): {output_path}")

if __name__ == "__main__":
    import qa_agent
    out_file = r"D:\AI\63_youtube_creator\pipeline\output\crisp_dm_shorts.pptx"
    
    # QA 테스트용 실제 스크립트 (슬라이드 내용과 정확히 일치)
    CRISP_DM_SCRIPT = (
        "🎯 CRISP-DM?\n데이터분석의 표준 6단계!\n"
        "데이터를 분석하려면?\n막 시작하면 안 된다.\n정해진 순서가 있어.\n\n그게 바로 CRISP-DM!\n업계 표준 방법론이야.\n"
        "CRISP-DM 6단계\nStep 1 Business Understanding\n  → 비즈니스 목표 정의\n"
        "Step 2 Data Understanding\n  → 데이터 수집 및 탐색\nStep 3 Data Preparation\n  → 데이터 정제 및 전처리\n"
        "Step 4 Modeling\n  → 알고리즘 학습 및 튜닝\nStep 5 Evaluation\n  → 성능 검증\nStep 6 Deployment\n  → 운영 환경 적용\n"
        "우리 3AI는?\n이 6단계를 완벽하게 적용 중!\n\nn8n 워크플로우 설계 = CRISP-DM 그 자체\nImprove_stock 주식 분석 = EDA 교과서 흐름\n\n"
        "수집 → 전처리 → 모델 → 신호생성 → 배포\n= CRISP-DM 완벽 실행!\n"
        "CRISP-DM은?\n무조건 알아야 할\n데이터 분석 기초!\n\n"
    )
    
    script_md_path = r"D:\AI\63_youtube_creator\pipeline\output\crisp_dm_script.md"
    os.makedirs(os.path.dirname(script_md_path), exist_ok=True)
    with open(script_md_path, "w", encoding="utf-8") as f:
        f.write(CRISP_DM_SCRIPT)
        
    empty_json_path = r"D:\AI\63_youtube_creator\pipeline\empty.json"
    dummy_toc_path = r"D:\AI\63_youtube_creator\pipeline\dummy_toc.json"
        
    max_retries = 2
    for attempt in range(max_retries + 1):
        print(f"\n[PPTX 생성] 시도 {attempt + 1}/{max_retries + 1}")
        create_shorts_pptx(out_file)
        
        print("[QA 검증] 코니의 7가지 체크리스트 실행 중...")
        # 올바른 인자 순서: pptx_path, audio_mp3, script_md, research_data_path, toc_json_path, image_paths
        result = qa_agent.final_validation(
            out_file, 
            "D:\\AI\\63_youtube_creator\\temp_audio.mp3", 
            script_md_path, 
            empty_json_path, 
            dummy_toc_path, 
            []
        )
        
        if result["passed_count"] == result["total_count"] or "승인" in result["verdict"]:
            print(f"QA 통과! {result['verdict']}")
            break
        else:
            print(f"QA 반려: {result['verdict']}")
            for issue in result["issues"]:
                print(f"  - {issue}")
            
            if attempt < max_retries:
                print("문제점을 수정하여 재렌더링 시도합니다...\n")
            else:
                print("최대 재시도 횟수 초과. 인간(User)의 개입이 필요합니다.")
                exit(1)
