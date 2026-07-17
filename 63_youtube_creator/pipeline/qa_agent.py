import os
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
from difflib import SequenceMatcher
import json

try:
    from moviepy.editor import AudioFileClip
except ImportError:
    from moviepy import AudioFileClip

def check_text_overflow(slide):
    issues = []
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            # Note: bounding_box is not always directly available on text_frame in python-pptx.
            # We approximate by shape bounding box.
            right_edge = shape.left + shape.width
            # Slide width for 1080x1920 (11.25 inches wide)
            # 1 inch = 914400 EMUs
            slide_width_emu = 11.25 * 914400
            
            # If right edge exceeds safe margin (e.g. 0.5 inches from right edge)
            safe_margin = slide_width_emu - int(0.5 * 914400)
            if right_edge > safe_margin:
                name = shape.name if shape.name else f"Shape_{i}"
                issues.append(f"[{name}] 텍스트 오버플로우 (안전영역 초과)")
                
    return len(issues) == 0, issues

def check_font_consistency(slides):
    title_sizes = []
    body_sizes = []
    
    for slide in slides:
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                if not shape.text_frame.paragraphs:
                    continue
                font = shape.text_frame.paragraphs[0].font
                if not font.size:
                    continue
                    
                name = shape.name.lower() if shape.name else ""
                
                # We identify role by shape name or roughly by size if names aren't set
                if "title" in name:
                    title_sizes.append((i, font.size))
                elif "body" in name or "text" in name:
                    body_sizes.append((i, font.size))

    title_std = Pt(44)
    body_std = Pt(24)
    
    issues = []
    for idx, size in title_sizes:
        if abs(size - title_std) > Pt(2):
            issues.append(f"제목 폰트 크기 불일치: {size.pt}pt (표준 44pt)")
            
    for idx, size in body_sizes:
        if abs(size - body_std) > Pt(1):
            issues.append(f"본문 폰트 크기 불일치: {size.pt}pt (표준 24pt)")
            
    return len(issues) == 0, issues

def check_image_specs(slide, image_paths):
    issues = []
    for img_path in image_paths:
        if not os.path.exists(img_path):
            continue
        try:
            img = Image.open(img_path)
            width, height = img.size
            
            if width < 1080 or height < 1080:
                issues.append(f"이미지 해상도 낮음 ({width}x{height})")
                
            aspect_ratio = width / height if height else 1
            if not (1.75 < aspect_ratio < 1.81 or 1.30 < aspect_ratio < 1.36 or 0.5 < aspect_ratio < 0.6):
                # Also allow 9:16 (~0.56) for vertical shorts backgrounds
                issues.append(f"비표준 이미지 비율 ({aspect_ratio:.2f})")
        except Exception as e:
            issues.append(f"이미지 검증 실패: {e}")
            
    return len(issues) == 0, issues

def check_content_consistency(script_text, slide_texts):
    issues = []
    for i, text in enumerate(slide_texts):
        if not text.strip():
            continue
        # Check if the slide text appears anywhere in the script or has high similarity
        matcher = SequenceMatcher(None, text.strip(), script_text)
        # Using a relaxed similarity for shorts since texts are chunked
        # Or check if text is a substring
        if text.strip() not in script_text:
            if not script_text.strip():
                similarity = 0
            else:
                similarity = max([SequenceMatcher(None, text.strip(), script_chunk).ratio() for script_chunk in script_text.split('\n') if script_chunk.strip()])
            
            if similarity < 0.7:  # Adjusted for practical use
                issues.append(f"슬라이드 {i+1}: 스크립트 불일치 의심 (일치도 {similarity*100:.1f}%)")
    
    return len(issues) == 0, issues

def check_hallucinations(slide_texts, research_facts):
    # Dummy implementation as per Coni's logic
    issues = []
    numbers = research_facts.get("numbers", [])
    
    for i, text in enumerate(slide_texts):
        words = text.replace(',', '').replace('.', ' ').split()
        for word in words:
            if word.isdigit():
                val = int(word)
                # If there are numbers in research_facts, ensure this number is valid
                if numbers and val not in numbers:
                    issues.append(f"슬라이드 {i+1}: 미확인 숫자 발견 ({val})")
                    
    return len(issues) == 0, issues

def get_luminance(rgb):
    r, g, b = [x/255.0 for x in rgb]
    r = r/12.92 if r <= 0.03928 else ((r+0.055)/1.055)**2.4
    g = g/12.92 if g <= 0.03928 else ((g+0.055)/1.055)**2.4
    b = b/12.92 if b <= 0.03928 else ((b+0.055)/1.055)**2.4
    return 0.2126*r + 0.7152*g + 0.0722*b

def check_color_contrast(slide):
    issues = []
    # Simplified contrast check. In PPTX extracting accurate bg color vs text color is complex 
    # without rendering. We'll do a placeholder check that assumes Dark BG for now.
    bg_rgb = (18, 18, 18) # 0x12, 0x12, 0x12 default shorts bg
    bg_lum = get_luminance(bg_rgb)
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.color and run.font.color.type == 1: # RGB type
                        r, g, b = run.font.color.rgb
                        txt_lum = get_luminance((r, g, b))
                        ratio = (max(bg_lum, txt_lum) + 0.05) / (min(bg_lum, txt_lum) + 0.05)
                        
                        required = 3.0
                        if ratio < required:
                            issues.append(f"명도 대비 부족: {ratio:.1f}:1 (요구 {required}:1)")
    
    return len(issues) == 0, issues

def check_timing_sync(audio_path, toc_json):
    issues = []
    try:
        audio = AudioFileClip(audio_path)
        audio_duration = audio.duration
        audio.close()
        
        toc_timings = [scene.get("duration_sec", 0) for scene in toc_json.get("scenes", [])]
        total_toc_time = sum(toc_timings)
        
        diff = abs(audio_duration - total_toc_time)
        if diff > len(toc_timings) * 0.5: # Tolerance: 0.5s per slide cumulative
            issues.append(f"전체 타이밍 오차 발생: 오디오({audio_duration:.1f}초) vs 대본({total_toc_time:.1f}초)")
            
    except Exception as e:
        issues.append(f"타이밍 검증 실패: {e}")
        
    return len(issues) == 0, issues

def final_validation(pptx_path, audio_mp3, script_md, research_data_path, toc_json_path, image_paths):
    all_passed = True
    all_issues = []
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"verdict": f"❌ 반려 (PPTX 로드 실패: {e})", "issues": [str(e)], "passed_count": 0, "total_count": 7}
        
    script_text = ""
    if os.path.exists(script_md):
        with open(script_md, "r", encoding="utf-8") as f:
            script_text = f.read()
            
    research_data = {}
    if os.path.exists(research_data_path):
        with open(research_data_path, "r", encoding="utf-8") as f:
            research_data = json.load(f)
            
    toc_json = {}
    if os.path.exists(toc_json_path):
        with open(toc_json_path, "r", encoding="utf-8") as f:
            toc_json = json.load(f)
    
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
        slide_texts.append("\n".join(texts))
        
        # 1-1
        p, i = check_text_overflow(slide)
        all_passed &= p
        all_issues.extend(i)
        
        # 1-3 (Simplified image spec check)
        p, i = check_image_specs(slide, image_paths)
        all_passed &= p
        all_issues.extend(i)
        
        # 2-3
        p, i = check_color_contrast(slide)
        all_passed &= p
        all_issues.extend(i)
        
    # 1-2
    p, i = check_font_consistency(prs.slides)
    all_passed &= p
    all_issues.extend(i)
    
    # 2-1
    p, i = check_content_consistency(script_text, slide_texts)
    all_passed &= p
    all_issues.extend(i)
    
    # 2-2
    p, i = check_hallucinations(slide_texts, research_data)
    all_passed &= p
    all_issues.extend(i)
    
    # 3-1
    p, i = check_timing_sync(audio_mp3, toc_json)
    all_passed &= p
    all_issues.extend(i)
    
    passed_count = 7 - len(all_issues)
    if all_passed and len(all_issues) == 0:
        verdict = "✅ 승인 (모든 항목 통과)"
    else:
        verdict = f"❌ 반려 ({len(all_issues)}개 문제 발생)"
        
    return {
        "verdict": verdict,
        "issues": all_issues,
        "passed_count": passed_count,
        "total_count": 7
    }

if __name__ == "__main__":
    # Test stub
    pass
